from pptx import Presentation
from pptx.exc import PackageNotFoundError
from typing import List, Dict, Optional

class PPTXProcessor:
    def __init__(self, presentation: Presentation):
        self.presentation = presentation

    @staticmethod
    def _extract_text_from_shape(shape) -> Optional[str]:
        """Extracts text from a single shape."""
        if not shape.has_text_frame:
            return None
        return shape.text

    @staticmethod
    def _extract_text_from_table(shape) -> Optional[List[str]]:
        """Extracts text from a table shape."""
        if not shape.has_table:
            return None
        return [cell.text for row in shape.table.rows for cell in row.cells if cell.text]

    @staticmethod
    def _extract_text_from_smartart(shape) -> Optional[List[str]]:
        """Extracts text from a SmartArt shape."""
        if shape.shape_type != 14:  # Shape type 14 corresponds to SmartArt
            return None

        smartart_text = []
        for i in range(shape.graphic_frame.graphic.graphic_data.element.child_len()):
            node = shape.graphic_frame.graphic.graphic_data.element.getchildren()[i]
            for paragraph in node.xpath('.//a:t', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                smartart_text.append(paragraph.text)
        return smartart_text

    def _extract_text_from_slide(self, slide) -> Dict[str, str]:
        """Extracts text, tables, and SmartArt from a slide."""
        slide_data = {"title": "", "subtitle": "", "body": "", "tables": [], "smartart": []}

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                slide_data["title"] = self._extract_text_from_shape(shape)
            elif len(slide.placeholders) > 1 and shape == slide.placeholders[1]:
                slide_data["subtitle"] = self._extract_text_from_shape(shape)
            else:
                text = self._extract_text_from_shape(shape)
                if text:
                    slide_data["body"] += text + "\n"

                table = self._extract_text_from_table(shape)
                if table:
                    slide_data["tables"].append(table)

                smartart = self._extract_text_from_smartart(shape)
                if smartart:
                    slide_data["smartart"].append(smartart)

        return slide_data

    def extract_text(self) -> List[Dict[str, str]]:
        """Extracts text, tables, and SmartArt from all slides in the presentation."""
        return [self._extract_text_from_slide(slide) for slide in self.presentation.slides]

def load_presentation(file_path: str) -> Presentation:
    """Loads a PowerPoint presentation from a file."""
    try:
        return Presentation(file_path)
    except PackageNotFoundError:
        raise ValueError(f"File not found at the path: {file_path}")

def extract_text_from_presentation(file_path: str) -> List[Dict[str, str]]:
    """Extracts text, tables, and SmartArt from a PowerPoint presentation."""
    presentation = load_presentation(file_path)
    processor = PPTXProcessor(presentation)
    return processor.extract_text()