from pptx import Presentation
from pptx.exc import PackageNotFoundError
from typing import List, Dict, Optional


class PPTXProcessor:
    def __init__(self, presentation: Presentation):
        self.presentation = presentation

    @staticmethod
    def _extract_text_from_slide(slide) -> Dict[str, str]:
        slide_data = {"title": "", "subtitle": "", "body": ""}

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    slide_data["title"] = shape.text
                elif len(slide.placeholders) > 1 and shape == slide.placeholders[
                    1]:  # Check for existence of placeholder 1
                    slide_data["subtitle"] = shape.text
                else:
                    slide_data["body"] += shape.text + "\n"

        return slide_data

    def extract_text(self) -> List[Dict[str, str]]:
        return [self._extract_text_from_slide(slide) for slide in self.presentation.slides]


def load_presentation(file_path: str) -> Presentation:
    try:
        return Presentation(file_path)
    except PackageNotFoundError:
        raise ValueError(f"File not found at the path: {file_path}")


def extract_text_from_presentation(file_path: str) -> List[Dict[str, str]]:
    """Extracts text from a PowerPoint presentation.

    Args:
        file_path: The path to the PowerPoint file.

    Returns:
        A list of dictionaries containing the title, subtitle, and body text from each slide.
    """
    presentation = load_presentation(file_path)
    processor = PPTXProcessor(presentation)
    return processor.extract_text()
